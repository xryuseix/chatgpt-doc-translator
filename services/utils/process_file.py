import os
from fastapi import HTTPException
from io import BufferedReader
from typing import Optional
from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import hashlib
import docx2txt
import csv
import pptx
from tqdm import tqdm

from models.models import Document


async def get_document_from_file(file: UploadFile) -> Document:
    extracted_text = await extract_text_from_form_file(file)
    doc = Document(text=extracted_text)

    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""
    print("Extracting text from file: ", filepath)

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    # Open the file in binary mode
    file = open(filepath, "rb")
    extracted_text = extract_text_from_file(file, mimetype)

    return extracted_text


def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = ""
        for page in tqdm(reader.pages):
            extracted_text += page.extract_text()
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in tqdm(reader):
            extracted_text += " ".join(row) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in tqdm(presentation.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        file.close()
        raise ValueError("Unsupported file type: {}".format(mimetype))

    file.close()
    return extracted_text


# Extract text from a file based on its mimetype
async def extract_text_from_form_file(file: UploadFile):
    """Return the text content of a file."""
    # get the file body from the upload file object
    mimetype = file.content_type
    print(f"mimetype: {mimetype}")

    file_stream = await file.read()

    hash_code = hashlib.sha256(file_stream).hexdigest()
    
    if not os.path.exists("./temp_files/"):
        os.makedirs("./temp_files/")
        print("Temporary Folder created successfully!")
        
    temp_file_path = f"./temp_files/{hash_code}"
    
    try:
        with open(temp_file_path, "wb") as f:
            f.write(file_stream)
        extracted_text = extract_text_from_filepath(temp_file_path, mimetype)
        
    except Exception as e:
        raise e

    os.remove(temp_file_path)

    return extracted_text
